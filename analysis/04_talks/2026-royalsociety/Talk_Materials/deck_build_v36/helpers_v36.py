"""
Shared helpers for v3.6 REVISED parallel slide builders.

Each slide builder file deck_build_v36/slides_v36/sNN.py imports from this
module and exposes a single function build_sNN(prs) that adds one slide to prs.

Universal fonts only on native PPT text:
    Calibri (body) · Georgia (headers, take-home) · Consolas (kicker, mono)

NO Crimson Pro / IBM Plex in native text.
"""
import os, subprocess
from pathlib import Path
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── Paths ──────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent   # Talk_Materials/
ASSETS = ROOT / "deck_assets"
ASSETS_V36 = ROOT / "deck_assets_v36"
FIGS = ROOT / "figs"
CLIPS = Path.home() / "Desktop" / "Herring_Talk_Clips" / "clips"
CLIPS_TOP = Path.home() / "Desktop" / "Herring_Talk_Clips"
POSTERS = CLIPS_TOP / "poster_frames"

# ─── Canvas ─────────────────────────────────────────────────────────────────
DECK_W_IN = 13.333
DECK_H_IN = 7.5

# ─── Palette ────────────────────────────────────────────────────────────────
C_DARK     = RGBColor(0x0E, 0x0E, 0x0E)
C_LIGHT    = RGBColor(0xFB, 0xFA, 0xF7)
C_INK      = RGBColor(0xF0, 0xEE, 0xE9)
C_INK_DARK = RGBColor(0x1C, 0x19, 0x16)
C_RUST     = RGBColor(0xD9, 0x71, 0x4F)
C_AMBER    = RGBColor(0xCF, 0xA0, 0x55)
C_MARINE   = RGBColor(0x6E, 0x9B, 0xC4)
C_KELP     = RGBColor(0x8A, 0xA0, 0x74)
C_SOFT     = RGBColor(0xA8, 0xA5, 0x9F)
C_SOFT_D   = RGBColor(0x6F, 0x68, 0x60)
C_RULE     = RGBColor(0x2E, 0x2C, 0x28)

# Universal fonts
HEAD = "Georgia"
BODY = "Calibri"
MONO = "Consolas"

# ─── Slide / shape helpers ──────────────────────────────────────────────────
def add_slide(prs, dark=False):
    BLANK = prs.slide_layouts[6]
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK if dark else C_LIGHT
    return s

def add_text(slide, text, x, y, w, h, *, font=BODY, size=18, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
             dark=False):
    if color is None:
        color = C_INK if dark else C_INK_DARK
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tx

def add_multi_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   line_spacing=1.15, dark=False):
    """runs: list of dicts {text, size?, bold?, italic?, font?, color?, new_para?}"""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, run in enumerate(runs):
        if i > 0 and run.get("new_para"):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = run["text"]
        f = r.font
        f.name = run.get("font", BODY)
        f.size = Pt(run.get("size", 18))
        f.bold = run.get("bold", False)
        f.italic = run.get("italic", False)
        color = run.get("color")
        if color is None:
            color = C_INK if dark else C_INK_DARK
        f.color.rgb = color
    return tx

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    sh.shadow.inherit = False
    return sh

def add_amber_rule(slide, x, y, w):
    return add_rect(slide, x, y, w, 0.04, C_AMBER)

def add_image(slide, path, x, y, w=None, h=None, send_to_back=False):
    path = str(path)
    if w is None and h is None:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y))
    elif w is None:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
    elif h is None:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    else:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if send_to_back:
        spTree = pic._element.getparent()
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)
    return pic

def add_video(slide, video_path, poster_path, x, y, w, h):
    """Embed a movie. Poster image is the on-slide thumbnail. ALL EMBEDDED."""
    return slide.shapes.add_movie(
        str(video_path),
        Inches(x), Inches(y), Inches(w), Inches(h),
        poster_frame_image=str(poster_path),
        mime_type="video/mp4",
    )

def kicker(slide, text, dark=False, x=0.4, y=0.3):
    return add_text(slide, text.upper(), x, y, 12.5, 0.4,
                    font=MONO, size=11, color=C_RUST,
                    align=PP_ALIGN.LEFT, dark=dark)

def title_h1(slide, text, dark=False, x=0.4, y=0.7, w=12.5, h=1.0, size=34):
    return add_text(slide, text, x, y, w, h,
                    font=HEAD, size=size, bold=True,
                    color=C_INK if dark else C_INK_DARK,
                    align=PP_ALIGN.LEFT, dark=dark)

def credit(slide, text, dark=False, x=0.4, y=7.1, w=12.5):
    return add_text(slide, text, x, y, w, 0.3,
                    font=MONO, size=8.5, color=C_SOFT if dark else C_SOFT_D,
                    align=PP_ALIGN.RIGHT, dark=dark)

def footage_credit(slide, dark=True, text="Footage: Pacific Wild"):
    """Pacific Wild credit, REQUIRED on every video slide."""
    return add_text(slide, text, DECK_W_IN - 3.0, DECK_H_IN - 0.4, 2.7, 0.3,
                    font=MONO, size=8.5, color=C_SOFT if dark else C_SOFT_D,
                    align=PP_ALIGN.RIGHT, dark=dark)

def speaker_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ─── Asset resolution ───────────────────────────────────────────────────────
def asset(*parts):
    """Look up an asset in deck_assets_v36, deck_assets, or figs (first hit)."""
    for base in (ASSETS_V36, ASSETS, FIGS):
        candidate = base.joinpath(*parts)
        if candidate.exists():
            return str(candidate)
    raise FileNotFoundError(f"Asset not found: {parts}")

def asset_opt(*parts):
    """Like asset() but returns None instead of raising."""
    try:
        return asset(*parts)
    except FileNotFoundError:
        return None

def clip(name):
    p = CLIPS / name
    if p.exists():
        return str(p)
    p = CLIPS_TOP / name
    if p.exists():
        return str(p)
    raise FileNotFoundError(f"Clip not found: {name}")

def poster(name):
    p = POSTERS / name
    if p.exists():
        return str(p)
    raise FileNotFoundError(f"Poster not found: {name}")
