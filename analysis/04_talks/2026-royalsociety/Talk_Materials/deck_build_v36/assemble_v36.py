#!/usr/bin/env python3
"""
Phase C assembler for the v3.6 REVISED Royal Society herring deck.

Imports every slide builder in spine order and writes the combined .pptx.
If a single builder fails, inserts a BROKEN placeholder slide and continues
so the rest of the deck still ships.
"""
import sys
import importlib
import traceback
from pathlib import Path

HERE = Path(__file__).resolve().parent
SLIDES_DIR = HERE / "slides_v36"

# Allow `from helpers_v36 import *` from inside slide modules
sys.path.insert(0, str(HERE))
# Allow `from sNN import build_sNN`
sys.path.insert(0, str(SLIDES_DIR))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = (
    HERE.parent
    / "Herring_RoyalSociety_Stier_2026_v3.6_REVISED.pptx"
)

SPINE = (
    [f"s{n:02d}" for n in range(1, 31)]
    + [f"x{n}" for n in range(1, 5)]
    + [f"b{n}" for n in range(1, 6)]
)


def add_broken_placeholder(prs, slide_id, err_text):
    """Insert a stark light slide explaining which builder blew up."""
    BLANK = prs.slide_layouts[6]
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFB, 0xFA, 0xF7)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"BROKEN SLIDE — {slide_id}"
    r.font.name = "Consolas"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xD9, 0x71, 0x4F)

    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = err_text
    r2.font.name = "Consolas"
    r2.font.size = Pt(11)
    r2.font.color.rgb = RGBColor(0x1C, 0x19, 0x16)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    patched = []
    failed = []

    for slide_id in SPINE:
        try:
            mod = importlib.import_module(slide_id)
            build_fn = getattr(mod, f"build_{slide_id}")
            build_fn(prs)
            print(f"[OK]   {slide_id}")
        except Exception as e:
            tb = traceback.format_exc()
            print(f"[FAIL] {slide_id}: {e}", file=sys.stderr)
            print(tb, file=sys.stderr)
            failed.append((slide_id, tb))
            add_broken_placeholder(prs, slide_id, tb)

    prs.save(str(OUT))
    print(f"\nSAVED: {OUT}")
    print(f"Total slides: {len(prs.slides)}")
    if failed:
        print(f"FAILED builders ({len(failed)}):")
        for sid, _ in failed:
            print(f"  - {sid}")
    return failed


if __name__ == "__main__":
    failed = main()
    sys.exit(0)
