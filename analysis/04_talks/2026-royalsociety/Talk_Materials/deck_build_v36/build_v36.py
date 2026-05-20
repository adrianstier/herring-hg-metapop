"""
v3.6 deck builder — Royal Society Herring talk.

Strategy: NATIVE python-pptx with UNIVERSAL fonts (Calibri / Georgia / Consolas)
so the deck is portable across machines (Royal Society PC, jump-drive). All
images and videos EMBEDDED in the .pptx (no linked references).

Output: ../Herring_RoyalSociety_Stier_2026_v3.6_REDESIGN.pptx
Fallback (untouched): ../Herring_RoyalSociety_Stier_2026_claude-code_20260519-152552_A-minus-4_SHIPPABLE.pptx
"""
import os, sys, subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ─── Paths ──────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent  # Talk_Materials
ASSETS = ROOT / "deck_assets"
ASSETS_V36 = ROOT / "deck_assets_v36"
FIGS = ROOT / "figs"
CLIPS = Path.home() / "Desktop" / "Herring_Talk_Clips" / "clips"
CLIPS_TOP = Path.home() / "Desktop" / "Herring_Talk_Clips"
POSTERS = CLIPS_TOP / "poster_frames"

OUT_PPTX = ROOT / "Herring_RoyalSociety_Stier_2026_v3.6_REDESIGN.pptx"

# ─── Canvas ─────────────────────────────────────────────────────────────────
DECK_W_IN = 13.333  # widescreen
DECK_H_IN = 7.5

# ─── Palette (warm white default, dark variant for title/close) ─────────────
C_DARK     = RGBColor(0x0E, 0x0E, 0x0E)
C_LIGHT    = RGBColor(0xFB, 0xFA, 0xF7)
C_INK      = RGBColor(0xF0, 0xEE, 0xE9)   # ink on dark
C_INK_DARK = RGBColor(0x1C, 0x19, 0x16)   # ink on light
C_RUST     = RGBColor(0xD9, 0x71, 0x4F)
C_AMBER    = RGBColor(0xCF, 0xA0, 0x55)
C_MARINE   = RGBColor(0x6E, 0x9B, 0xC4)
C_KELP     = RGBColor(0x8A, 0xA0, 0x74)
C_SOFT     = RGBColor(0xA8, 0xA5, 0x9F)
C_SOFT_D   = RGBColor(0x6F, 0x68, 0x60)
C_RULE     = RGBColor(0x2E, 0x2C, 0x28)

# Universal fonts (present on Win/Mac Royal Society PC)
HEAD = "Georgia"     # serif headers (substitute for Crimson Pro)
BODY = "Calibri"     # body (substitute for IBM Plex Sans)
MONO = "Consolas"    # mono (substitute for IBM Plex Mono)

# ─── Build prs ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(DECK_W_IN)
prs.slide_height = Inches(DECK_H_IN)
BLANK = prs.slide_layouts[6]  # blank layout (no placeholders)

# ─── Helpers ────────────────────────────────────────────────────────────────
def add_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK if dark else C_LIGHT
    return s

def add_text(slide, text, x, y, w, h, *, font=BODY, size=18, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
             char_spacing=None, dark=False):
    if color is None:
        color = C_INK if dark else C_INK_DARK
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tx

def add_multi_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   line_spacing=1.15, dark=False):
    """runs is list of dicts: {text, size, bold, italic, font, color, break_after}"""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, run in enumerate(runs):
        if i > 0 and run.get("new_para"):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = run["text"]
        f = r.font
        f.name = run.get("font", BODY)
        f.size = Pt(run.get("size", 18))
        f.bold = run.get("bold", False)
        f.italic = run.get("italic", False)
        color = run.get("color")
        if color is None:
            color = C_INK if dark else C_INK_DARK
        f.color.rgb = color
    return tx

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    sh.shadow.inherit = False
    return sh

def add_amber_rule(slide, x, y, w):
    return add_rect(slide, x, y, w, 0.04, C_AMBER)

def add_image(slide, path, x, y, w=None, h=None):
    if w is None and h is None:
        return slide.shapes.add_picture(path, Inches(x), Inches(y))
    if w is None:
        return slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
    if h is None:
        return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))

def add_video(slide, video_path, poster_path, x, y, w, h):
    """Embed a video. Uses poster image as the on-slide thumbnail."""
    # python-pptx supports add_movie with poster_frame_image
    mov = slide.shapes.add_movie(
        video_path,
        Inches(x), Inches(y),
        Inches(w), Inches(h),
        poster_frame_image=poster_path,
        mime_type="video/mp4"
    )
    return mov

def kicker(slide, text, dark=False, x=0.4, y=0.3):
    return add_text(slide, text.upper(), x, y, 12.5, 0.4,
                    font=MONO, size=11, color=C_RUST,
                    align=PP_ALIGN.LEFT, char_spacing=2, dark=dark)

def title_h1(slide, text, dark=False, x=0.4, y=0.7, w=12.5, h=1.0, size=34):
    return add_text(slide, text, x, y, w, h,
                    font=HEAD, size=size, bold=True,
                    color=C_INK if dark else C_INK_DARK,
                    align=PP_ALIGN.LEFT, dark=dark)

def credit(slide, text, dark=False, x=0.4, y=7.1, w=12.5):
    return add_text(slide, text, x, y, w, 0.3,
                    font=MONO, size=8.5, color=C_SOFT if dark else C_SOFT_D,
                    align=PP_ALIGN.RIGHT, dark=dark)

def speaker_note(slide, text):
    """Add speaker notes for the presenter view."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = text

# ─── Asset resolution helpers ───────────────────────────────────────────────
def asset(*parts):
    """Look up an asset in deck_assets or deck_assets_v36."""
    for base in (ASSETS_V36, ASSETS, FIGS):
        candidate = base.joinpath(*parts)
        if candidate.exists():
            return str(candidate)
    raise FileNotFoundError(f"Asset not found: {parts}")

def clip(name):
    p = CLIPS / name
    if p.exists():
        return str(p)
    # try top-level
    p = CLIPS_TOP / name
    if p.exists():
        return str(p)
    raise FileNotFoundError(f"Clip not found: {name}")

def poster(name):
    p = POSTERS / name
    if p.exists():
        return str(p)
    raise FileNotFoundError(f"Poster not found: {name}")

# ─── Frame extraction (for S3 contrast: pull pre-plume frame from v2_01) ────
def ensure_extracted_frame(video_path, time_sec, out_path):
    """Extract a frame from a video at time_sec via ffmpeg if not already done."""
    if Path(out_path).exists():
        return out_path
    cmd = [
        "ffmpeg", "-y", "-ss", str(time_sec), "-i", video_path,
        "-vframes", "1", "-q:v", "2", out_path
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    return out_path

# ─── Slide builders ─────────────────────────────────────────────────────────

def s01_title():
    """S1 — Title + Thesis + Take-home"""
    s = add_slide(dark=True)
    # background image
    try:
        img = asset("01_title.png")
        pic = add_image(s, img, 0, 0, DECK_W_IN, DECK_H_IN)
        # send to back
        spTree = pic._element.getparent()
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)
    except FileNotFoundError:
        pass
    # dark scrim at bottom for type legibility
    add_rect(s, 0, 3.0, DECK_W_IN, DECK_H_IN - 3.0, C_DARK)
    # actually want partial-opacity scrim — replace with semi-transparent solid
    # python-pptx doesn't directly support alpha on solid fill; use shape transparency via XML
    # for portability, keep solid dark; image already dark
    # kicker
    kicker(s, "Session 5 · Tipping Points in Ecosystem Services", dark=True, y=0.3)
    # h1
    title_h1(s, "Pacific Herring at Haida Gwaii", dark=True,
             y=3.4, size=48)
    # subtitle
    add_text(s, "A coupled social–ecological system — drawn down twice in one lifetime, and back only once.",
             0.4, 4.6, 12.5, 0.7, font=HEAD, size=20, italic=True,
             color=C_SOFT, align=PP_ALIGN.LEFT, dark=True)
    # take-home with amber rule
    add_amber_rule(s, 0.4, 5.7, 0.06)
    add_text(s, '"Recovery is a moving target."',
             0.6, 5.6, 12.5, 0.6, font=HEAD, size=24, italic=True, bold=True,
             color=C_AMBER, align=PP_ALIGN.LEFT, dark=True)
    # credit
    credit(s, "Adrian Stier · UC Santa Barbara · US–UK Forum · Session 5 · 20 May 2026", dark=True, y=7.0)
    # speaker note
    speaker_note(s, "Pacific herring at Haida Gwaii — one of the longest-studied "
                 "forage-fish systems on the planet — has collapsed twice in one "
                 "lifetime. The first time it came back in five years. The second "
                 "time it has not come back in thirty. Recovery is a moving target.")

def s02_haida_gwaii():
    """S2 — Haida Gwaii: place + tension"""
    s = add_slide(dark=True)
    # full-bleed dark coastline; reuse 02_shore.png or fall back to 01_title
    try:
        img = asset("02_shore.png")
    except FileNotFoundError:
        img = asset("01_title.png")
    pic = add_image(s, img, 0, 0, DECK_W_IN, DECK_H_IN)
    # scrim for type
    add_rect(s, 0, DECK_H_IN - 2.5, DECK_W_IN, 2.5, C_DARK)
    kicker(s, "The Place", dark=True, y=0.3)
    title_h1(s, "Haida Gwaii", dark=True, y=5.3, size=48)
    add_text(s, "An archipelago at the productive edge of the Pacific — ~80 km off the British Columbia coast.",
             0.4, 6.4, 12.5, 0.5, font=HEAD, size=18, italic=True,
             color=C_SOFT, align=PP_ALIGN.LEFT, dark=True)
    credit(s, "Adrian C. Stier · Royal Society Forum 2026", dark=True, y=7.05)
    speaker_note(s, "Haida Gwaii. An archipelago at the productive edge of the Pacific "
                 "— about eighty kilometres off the British Columbia coast. One of the "
                 "wildest, most productive stretches of coastline in North America. "
                 "People have been here for ten thousand years. And for most of the year, "
                 "you would never know what's about to happen.")

def s03_contrast():
    """S3 — Contrast: water without/with herring (side-by-side)"""
    s = add_slide(dark=False)
    # Try to extract paired frames from v2_01
    try:
        v201 = clip("v2_01_aerial_whitewater_spawn.mp4")
        before = ASSETS_V36 / "s03_before_water.png"
        after  = ASSETS_V36 / "s03_after_plume.png"
        ensure_extracted_frame(v201, 0.1, str(before))   # pre-plume (or close)
        ensure_extracted_frame(v201, 5.0, str(after))    # full plume
    except Exception as e:
        # fall back to repo assets
        before = asset("02_shore.png")
        after  = poster("v2_01_aerial_whitewater_spawn.jpg")

    kicker(s, "The Signature", y=0.3)
    # split panel: two equal halves
    half_w = (DECK_W_IN - 0.8) / 2
    panel_top = 1.2
    panel_h = DECK_H_IN - 2.6
    # left
    pl = add_image(s, str(before), 0.4, panel_top, w=half_w, h=panel_h)
    # right
    pr = add_image(s, str(after), 0.4 + half_w + 0.0, panel_top, w=half_w, h=panel_h)
    # labels
    add_text(s, "MOST OF THE YEAR", 0.5, panel_top + 0.1, half_w - 0.2, 0.3,
             font=MONO, size=12, bold=True, color=C_AMBER, align=PP_ALIGN.LEFT)
    add_text(s, "WHEN THE HERRING ARRIVE", 0.5 + half_w, panel_top + 0.1, half_w - 0.2, 0.3,
             font=MONO, size=12, bold=True, color=C_AMBER, align=PP_ALIGN.LEFT)
    # bottom caption strip
    add_amber_rule(s, 0.4, panel_top + panel_h + 0.15, 0.06)
    add_text(s, '"In a matter of days."',
             0.6, panel_top + panel_h + 0.1, 12.5, 0.45,
             font=HEAD, size=20, italic=True, color=C_INK_DARK, align=PP_ALIGN.LEFT)
    credit(s, "Footage: Pacific Wild", y=7.1)
    speaker_note(s, "This is what the water looks like. On the left — most of the year. "
                 "Dark, clear, cold. On the right — when the herring arrive. The whole "
                 "bay turns turquoise. You can see it from a satellite. The shift happens "
                 "in days.")

def s04_herring_arrive():
    """S4 — When the herring arrive: aerial → underwater (embedded video)"""
    s = add_slide(dark=True)
    # Embed v2_01 aerial spawn — autoplay-style
    try:
        v_path = clip("v2_01_aerial_whitewater_spawn.mp4")
        p_path = poster("v2_01_aerial_whitewater_spawn.jpg")
        # full-bleed video — left half aerial, right half underwater
        # For simplicity in a single slide, use ONE video full-bleed (the aerial),
        # and we'll dedicate next slide to the underwater school
        add_video(s, v_path, p_path, 0, 0, DECK_W_IN, DECK_H_IN)
    except Exception as e:
        # fallback: use poster image
        try:
            img = poster("v2_01_aerial_whitewater_spawn.jpg")
            add_image(s, img, 0, 0, DECK_W_IN, DECK_H_IN)
        except Exception:
            pass
    # minimal text — bottom-left subtle
    kicker(s, "", dark=True, y=0.3)
    credit(s, "Footage: Pacific Wild", dark=True, y=7.1)
    speaker_note(s, "This is the change you're seeing — made by these. "
                 "[Aerial spawn plume plays; cut to underwater on next slide if separate, "
                 "or contained here as a single full-bleed reveal.]")

def s04b_underwater():
    """S4 part 2 — underwater pure-herring schooling (the agents)"""
    s = add_slide(dark=True)
    try:
        v_path = str(CLIPS_TOP / "herring_schooling.mp4")
        if not Path(v_path).exists():
            v_path = clip("v2_02_underwater_school.mp4")
        # Use poster for v2_02 as fallback poster
        try:
            p_path = poster("v2_02_underwater_school.jpg")
        except FileNotFoundError:
            p_path = str(ASSETS_V36 / "s04b_poster.png")
            if not Path(p_path).exists():
                ensure_extracted_frame(v_path, 1.0, p_path)
        add_video(s, v_path, p_path, 0, 0, DECK_W_IN, DECK_H_IN)
    except Exception:
        try:
            img = poster("v2_02_underwater_school.jpg")
            add_image(s, img, 0, 0, DECK_W_IN, DECK_H_IN)
        except Exception:
            pass
    credit(s, "Footage: Pacific Wild", dark=True, y=7.1)
    speaker_note(s, '"— made by these." [Underwater verified pure-herring schooling plays.]')

def s05_wildlife():
    """S5 — And everything else comes with them"""
    s = add_slide(dark=True)
    try:
        v_path = clip("v2_04_bear_and_white_wolf.mp4")
        p_path = poster("v2_04_bear_and_white_wolf.jpg")
        add_video(s, v_path, p_path, 0, 0, DECK_W_IN, DECK_H_IN)
    except Exception:
        try:
            img = poster("v2_04_bear_and_white_wolf.jpg")
            add_image(s, img, 0, 0, DECK_W_IN, DECK_H_IN)
        except Exception:
            pass
    # text overlay bottom — IF visible against video
    add_rect(s, 0, DECK_H_IN - 1.4, DECK_W_IN, 1.4, C_DARK)
    add_text(s, "Bears. Wolves. Eagles. Sea lions. The whole coast eats.",
             0.4, 6.4, 12.5, 0.7, font=HEAD, size=24, italic=True, bold=True,
             color=C_INK, align=PP_ALIGN.LEFT, dark=True)
    credit(s, "Footage: Pacific Wild", dark=True, y=7.1)
    speaker_note(s, "And everything else comes with them. Bears. Wolves. Eagles. "
                 "Sea lions. The whole coast eats. [Quick cuts: bear+wolf → eagles "
                 "catching fish → sea lion underwater → sea lion colony.]")

def s06_predators():
    """S6 — The predators came back"""
    s = add_slide(dark=False)
    kicker(s, "What Else Changed", y=0.3)
    title_h1(s, "The predators came back.", y=0.7, size=36)
    # Hero image left
    try:
        img = poster("v1_15_sea_lion_colony.jpg")
        add_image(s, img, 0.4, 1.8, w=7.5, h=4.5)
    except Exception:
        pass
    # Data callouts right
    runs = [
        {"text": "Humpback whales\n", "size": 14, "bold": True, "font": MONO, "color": C_RUST},
        {"text": "~1,000 (1970s) → ", "size": 18, "font": BODY},
        {"text": "25,000–33,000", "size": 22, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": " (2021)\n\n", "size": 18, "font": BODY},
        {"text": "Harbour seals\n", "size": 14, "bold": True, "font": MONO, "color": C_RUST, "new_para": True},
        {"text": "~10,000 → ", "size": 18, "font": BODY},
        {"text": "~105,000", "size": 22, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "\n\n", "size": 18, "font": BODY},
        {"text": "Steller sea lions\n", "size": 14, "bold": True, "font": MONO, "color": C_RUST, "new_para": True},
        {"text": ">4× ", "size": 22, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "since 1970", "size": 18, "font": BODY},
    ]
    add_multi_text(s, runs, 8.2, 1.8, 4.7, 4.5)
    # punch line with amber rule
    add_amber_rule(s, 0.4, 6.45, 0.06)
    add_text(s, '"The mouth eating herring now is bigger than the fishery ever was."',
             0.6, 6.4, 12.5, 0.5, font=HEAD, size=18, italic=True,
             color=C_INK_DARK, align=PP_ALIGN.LEFT)
    credit(s, "Footage: Pacific Wild · Pop. data: DFO + literature synthesis", y=7.1)
    speaker_note(s, "But the coast you just saw — full of bears and wolves and eagles "
                 "and sea lions — looks like that for a different reason today than it "
                 "did fifty years ago. In the 1960s, commercial whaling had left the "
                 "ocean nearly predator-free. Since then: humpbacks went from about a "
                 "thousand to twenty-five to thirty-five thousand. Sea lions doubled, "
                 "then doubled again. The mouth eating herring today is bigger than "
                 "the fishery ever was.")

def s07_haida():
    """S7 — Introduce the Haida"""
    s = add_slide(dark=True)
    # Asset NEEDED — placeholder for Adrian to swap in a thriving-people image
    placeholder = ASSETS_V36 / "s07_haida_PLACEHOLDER.png"
    if not placeholder.exists():
        # create a simple dark placeholder with text "ASSET NEEDED"
        from PIL import Image, ImageDraw, ImageFont
        im = Image.new("RGB", (3840, 2160), (14, 14, 14))
        d = ImageDraw.Draw(im)
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Georgia.ttf", 60)
        except Exception:
            font = ImageFont.load_default()
        msg = ("[S7 ASSET NEEDED]\n\nReplace this image with a contemporary photograph\n"
               "of Haida people being Haida now — active, present-tense,\n"
               "agency-forward. (See narrative S7 entry for guidance.)\n\n"
               "Adrian provides — sourced with explicit consent.")
        d.multiline_text((600, 700), msg, fill=(208, 184, 144), font=font, spacing=20)
        im.save(placeholder)
    add_image(s, str(placeholder), 0, 0, DECK_W_IN, DECK_H_IN)
    # text overlay
    add_rect(s, 0, DECK_H_IN - 2.5, DECK_W_IN, 2.5, C_DARK)
    kicker(s, "The Haida", dark=True, y=0.3)
    title_h1(s, "On these islands for 10,000 years — and here now.",
             dark=True, y=5.4, size=32)
    credit(s, "Image: ASSET NEEDED — Adrian provides (Haida-source)", dark=True, y=7.1)
    speaker_note(s, "The Haida are an Indigenous nation. These islands are their home "
                 "— Haida Gwaii. They have been here for at least ten thousand years "
                 "and they are here now. Their language, ceremony, and marine knowledge "
                 "are alive and active. The Council of the Haida Nation is their "
                 "government, and since 1985 the Haida and the Canadian state have "
                 "co-governed parts of the archipelago together.")

def s08_archaeological():
    """S8 — Archaeological record: 10,000 years of herring at the centre"""
    s = add_slide(dark=False)
    kicker(s, "The Record", y=0.3)
    title_h1(s, "10,000 years of herring at the centre.", y=0.7, size=34)
    # Try McKechnie figure or baseline asset
    try:
        img = asset("04_baseline.png")
        add_image(s, img, 0.4, 1.8, w=8.0, h=4.6)
    except FileNotFoundError:
        pass
    # data callouts right column
    runs = [
        {"text": "49%", "size": 56, "bold": True, "font": HEAD, "color": C_RUST},
        {"text": " of fish bones\n", "size": 14, "font": BODY},
        {"text": "across 171 archaeological sites\n\n", "size": 12, "font": BODY, "color": C_SOFT_D},
        {"text": "<±10%", "size": 36, "bold": True, "font": HEAD, "color": C_INK_DARK, "new_para": True},
        {"text": " regional variance\n", "size": 14, "font": BODY},
        {"text": "over ~10,700 years\n\n", "size": 12, "font": BODY, "color": C_SOFT_D},
        {"text": "kʼaaw harvest evidence in middens\n", "size": 14, "italic": True, "font": HEAD, "new_para": True},
        {"text": "Language: iinang = 'plentiful' = herring", "size": 12, "italic": True, "font": HEAD, "color": C_SOFT_D},
    ]
    add_multi_text(s, runs, 8.8, 1.8, 4.1, 4.6)
    credit(s, "McKechnie et al. 2014 PNAS · HMTK", y=7.1)
    speaker_note(s, "And the record shows it. Across one hundred and seventy-one "
                 "archaeological sites — over ten thousand seven hundred years — "
                 "herring is roughly half of the fish bone we find. With less than ten "
                 "percent variance across the region. For ten millennia this fish was "
                 "the dominant fish people ate, and it was both variable and reliable "
                 "at the same time. The Haida even encoded it in their language — "
                 "iinang means 'plentiful' and it's the word for herring.")

def s09_three_keystones():
    """S9 — Three keystones, one fish (summary + transition)"""
    s = add_slide(dark=False)
    kicker(s, "The Setup", y=0.3)
    title_h1(s, "Three keystones, one fish.", y=0.7, size=34)
    # 3 equal columns
    col_w = (DECK_W_IN - 0.8 - 0.6) / 3  # 0.4 margins each side, 0.3 between
    cols = ["ECOLOGICAL", "CULTURAL", "ECONOMIC"]
    subhead = ["The channel", "10,000 years", "150 years"]
    body = [
        [("Bears · wolves · eagles · sea lions · salmon", BODY, 14, False),
         ("The wasp-waist of the coastal food web.", BODY, 12, True),
         ('"The whole coast eats."', HEAD, 16, True)],
        [("kʼaaw — first food of the spring.", BODY, 14, False),
         ("Trade currency · ceremony · language.", BODY, 12, True),
         ('"Part of the essence of who you are."', HEAD, 14, True),
         ("— Barbara Wilson, 2016", MONO, 10, False)],
        [("Reduction era → kazunoko luxury → present absence.", BODY, 14, False),
         ("1995 SOK peak $62.88/lb → $11–14/lb (2004).", BODY, 12, True),
         ("Last HG commercial roe fishery: 2002.", HEAD, 14, True),
         ("Zero since.", MONO, 11, False)],
    ]
    col_top = 2.0
    col_h = 4.2
    for i, (col, sh, b) in enumerate(zip(cols, subhead, body)):
        x = 0.4 + i * (col_w + 0.3)
        # column header
        add_text(s, col, x, col_top, col_w, 0.4,
                 font=MONO, size=14, bold=True, color=C_RUST, align=PP_ALIGN.LEFT)
        # subhead
        add_text(s, sh, x, col_top + 0.45, col_w, 0.4,
                 font=HEAD, size=18, bold=True, italic=True, color=C_INK_DARK)
        # body rows
        runs = []
        for j, (text, font, size, italic) in enumerate(b):
            if j > 0:
                runs.append({"text": "\n\n", "size": size, "font": font})
            runs.append({"text": text, "size": size, "font": font, "italic": italic})
        add_multi_text(s, runs, x, col_top + 1.0, col_w, col_h - 1.0)
    # bottom strip
    add_amber_rule(s, 0.4, 6.5, 0.06)
    add_text(s, '"And in one lifetime, the fish changed twice."',
             0.6, 6.45, 12.5, 0.5, font=HEAD, size=20, italic=True,
             color=C_INK_DARK, align=PP_ALIGN.LEFT)
    credit(s, "$ figures: 2024 HG Herring Rebuilding Plan §5.2.3", y=7.1)
    speaker_note(s, "So step back. Three keystone relationships on one fish. "
                 "Ecological — herring is the channel from ocean plankton to bears, "
                 "wolves, eagles, sea lions, salmon. The whole coast eats. "
                 "Cultural — for the Haida, ten thousand years of relationship; "
                 "kʼaaw is the first food of the spring. "
                 "Economic — a hundred and fifty years of commercial harvest layered "
                 "on top: the reduction era, then the kazunoko luxury market, then "
                 "absence. Three keystones, one fish. And in one lifetime, the fish "
                 "changed twice.")

def s10_heterogeneity():
    """S10 — Look closer: at the cove scale it's a different story"""
    s = add_slide(dark=False)
    kicker(s, "Look Closer", y=0.3)
    title_h1(s, "At the cove scale, it's a different story.", y=0.7, size=32)
    add_text(s, "Eleven subpopulations, one archipelago. Each cove ran on its own clock.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    # figure
    try:
        img = asset("08b_subpop_portfolio.png")
    except FileNotFoundError:
        try:
            img = asset("sb1_portfolio_periods.png")
        except FileNotFoundError:
            img = None
    if img:
        add_image(s, img, 0.6, 2.1, w=12.0, h=4.6)
    credit(s, "Data: DFO Pacific Herring assessments · Stier et al. 2020 Ecosphere", y=7.1)
    speaker_note(s, "Step closer. The same fish, the same archipelago — but at the cove "
                 "scale, almost every section was running on its own clock. Some coves "
                 "peaked when others crashed. Some crashed and recovered while others "
                 "stayed quiet for a decade. This is what the spawn-index data look like "
                 "when you stop averaging and start looking section by section. The whole "
                 "archipelago looked relatively steady because the parts were out of "
                 "phase. When parts are out of phase, the whole is buffered. That's a "
                 "portfolio — and we'll come back to that word in a minute.")

def s11_portfolio_theory():
    """S11 — Diversification removes variance, not yield (portfolio theory)"""
    s = add_slide(dark=False)
    kicker(s, "What Is a Portfolio?", y=0.3)
    title_h1(s, "Diversification removes variance, not yield.", y=0.7, size=32)
    add_text(s, "The resilience of wildly variable systems lives in their structure — not their mean.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    # try the animation poster, or fall back to async/sync illustration
    try:
        img = asset("sim_anim_v2_poster.jpg")
    except FileNotFoundError:
        try:
            img = asset("08b_subpop_portfolio.png")
        except FileNotFoundError:
            img = None
    if img:
        add_image(s, img, 0.6, 2.1, w=8.0, h=4.6)
    # right sidebar
    runs = [
        {"text": "Markowitz, 1952\n", "size": 14, "bold": True, "font": MONO, "color": C_RUST},
        {"text": "Diversification removes variance, not yield.\n\n", "size": 14, "font": BODY},
        {"text": "Same math in ecology\n", "size": 14, "bold": True, "font": MONO, "color": C_RUST, "new_para": True},
        {"text": "A metapopulation of asynchronous sub-stocks IS a portfolio.\n\n", "size": 14, "font": BODY},
        {"text": "Schindler et al., 2010 Nature\n", "size": 14, "bold": True, "font": MONO, "color": C_RUST, "new_para": True},
        {"text": "Demonstrated for Bristol Bay sockeye — the textbook ecological example.", "size": 14, "font": BODY},
    ]
    add_multi_text(s, runs, 8.9, 2.1, 4.0, 4.6)
    # data window callout
    add_text(s, "Haida Gwaii spawn-thickness · 1950–1967 · 11 sections",
             0.6, 6.55, 8.0, 0.3, font=MONO, size=10, color=C_SOFT_D, align=PP_ALIGN.LEFT)
    credit(s, "Markowitz 1952 · Schindler et al. 2010 · Doak et al. · Tilman · Loreau & de Mazancourt", y=7.1)
    speaker_note(s, "Quick detour into the theory. In 1952 an economist named Harry "
                 "Markowitz won a Nobel Prize for showing that if you hold a portfolio "
                 "of assets that don't all move together — that are imperfectly "
                 "correlated — you get the same long-run return as any one of them but "
                 "with much less variance. Diversification removes risk without removing "
                 "yield. In ecology, it's the same math. A metapopulation of asynchronous "
                 "sub-stocks IS a portfolio. From 1950 to the late 1960s, this is what "
                 "the Haida Gwaii data look like — eleven coves, each on its own clock, "
                 "and the regional whole stays steady. Schindler and colleagues showed "
                 "this for Bristol Bay sockeye in 2010. The structure does the work. "
                 "Asynchrony IS the buffer.")

def s12_portfolio_loss():
    """S12 — The portfolio came apart (synchrony rose at 1994)"""
    s = add_slide(dark=False)
    kicker(s, "The Portfolio Came Apart", y=0.3)
    title_h1(s, "The parts came into phase.", y=0.7, size=34)
    add_text(s, "Synchrony rose substantially post-1994; the buffer disappeared.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    # figure
    try:
        img = asset("09_synchrony.png")
    except FileNotFoundError:
        try:
            img = asset("sim_anim_v2_poster.jpg")
        except FileNotFoundError:
            img = None
    if img:
        add_image(s, img, 0.6, 2.1, w=8.0, h=4.6)
    # big data callout right
    runs = [
        {"text": "φ pre-1994\n", "size": 14, "font": MONO, "color": C_SOFT_D},
        {"text": "0.17\n\n", "size": 56, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "φ post-1994\n", "size": 14, "font": MONO, "color": C_SOFT_D, "new_para": True},
        {"text": "0.28\n\n", "size": 56, "bold": True, "font": HEAD, "color": C_RUST},
        {"text": ">60% increase\n", "size": 18, "bold": True, "font": HEAD, "color": C_RUST, "new_para": True},
        {"text": "Stier et al. 2020 Ecosphere", "size": 11, "italic": True, "font": MONO, "color": C_SOFT_D},
    ]
    add_multi_text(s, runs, 8.9, 2.1, 4.0, 4.6)
    credit(s, "Data: DFO · Synchrony metric: Stier et al. 2020 Ecosphere (audited per claim-control sheet)", y=7.1)
    speaker_note(s, "And then something changed. Around 1994, the cove-level cycles "
                 "started to come into phase with each other. Synchrony rose more than "
                 "sixty percent — what you're watching is the portfolio coming apart. "
                 "Asynchronous sub-stocks were buffering the whole; when they come into "
                 "phase, the buffer disappears. And here's the load-bearing finding: "
                 "total roe value held high through the 1980s while the structure that "
                 "produced it was already failing. Value lagged structure by a decade. "
                 "The catch statistics told the right story — too late.")

def s13_puzzle():
    """S13 — THE PUZZLE: two collapses, one lifetime, opposite recoveries"""
    s = add_slide(dark=False)
    kicker(s, "The Puzzle", y=0.3)
    title_h1(s, "Same fish. Same coast. Twice collapsed. Why did one come back?",
             y=0.7, size=28)
    # figure
    try:
        img = asset("05_two_collapses.png")
    except FileNotFoundError:
        try:
            img = asset("dfo_spawning_biomass.png")
        except FileNotFoundError:
            img = None
    if img:
        add_image(s, img, 0.6, 1.8, w=12.0, h=4.8)
    # annotations
    add_text(s, "1967: collapse → 5-yr recovery", 1.0, 6.5, 6.0, 0.4,
             font=MONO, size=12, bold=True, color=C_KELP, align=PP_ALIGN.LEFT)
    add_text(s, "Post-1994: collapse → no recovery in 30 yr", 7.0, 6.5, 6.0, 0.4,
             font=MONO, size=12, bold=True, color=C_RUST, align=PP_ALIGN.RIGHT)
    credit(s, "Source: DFO Pacific Herring SR 2024/2025", y=7.1)
    speaker_note(s, "Now back out to the aggregate. Two collapses in one lifetime. "
                 "The first — driven by industrial fishing — closed in 1967 and the "
                 "population came back in about five years. The second — beginning in "
                 "the 1990s — has not recovered in thirty. Same fish, same coast, "
                 "opposite outcomes. Why?")

def s14_how_we_know():
    """S14 — How we know: what we have to interrogate the puzzle"""
    s = add_slide(dark=False)
    kicker(s, "How We Know", y=0.3)
    title_h1(s, "Seventy-five years of data; ten thousand years of context; eleven subpopulations resolved.",
             y=0.7, size=22)
    # 4-card grid
    cards = [
        ("75 yr", "DFO spawn-thickness surveys\n1950–present · 11 sections"),
        ("10,000 yr", "Archaeological baseline\n(per S8 — McKechnie 2014)"),
        ("m1_stier_11", "Bayesian state-space\nmetapopulation model"),
        ("Audited", "Predator demand integrated as\nexternal pressure — never fitted"),
    ]
    card_w = (DECK_W_IN - 0.8 - 0.9) / 4  # 4 cards, 3 gaps of 0.3
    card_top = 2.5
    card_h = 3.8
    for i, (head, body) in enumerate(cards):
        x = 0.4 + i * (card_w + 0.3)
        # card background
        add_rect(s, x, card_top, card_w, card_h, RGBColor(0xF4, 0xF1, 0xEC))
        # head
        add_text(s, head, x + 0.2, card_top + 0.3, card_w - 0.4, 0.8,
                 font=HEAD, size=28, bold=True, color=C_RUST, align=PP_ALIGN.LEFT)
        # body
        add_text(s, body, x + 0.2, card_top + 1.3, card_w - 0.4, card_h - 1.6,
                 font=BODY, size=13, color=C_INK_DARK, align=PP_ALIGN.LEFT, line_spacing=1.3)
    credit(s, "DFO · McKechnie et al. 2014 · m1_stier_11 (Output/diagnostics/) · audited predator brief", y=7.1)
    speaker_note(s, "How do we interrogate this? Seventy-five years of DFO spawn-"
                 "thickness data, resolved at eleven sub-population sites. Ten thousand "
                 "years of archaeological baseline. A Bayesian state-space metapopulation "
                 "model — m1-stier-11 — across the sub-population network. Predator "
                 "demand enters as an audited external pressure, never as a fitted "
                 "parameter. So we can interrogate the suspects.")

def s15_ocean():
    """S15 — Suspect 1: ocean (necessary, not the barrier)"""
    s = add_slide(dark=False)
    kicker(s, "Suspect 1: Ocean", y=0.3)
    title_h1(s, "Necessary, not the barrier.", y=0.7, size=34)
    add_text(s, "Cool, productive years drove the 1960s rebound. Comparable cool years since 2000 did not.",
             0.4, 1.4, 12.5, 0.5, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    try:
        img = asset("06_climate_pdo.png")
        add_image(s, img, 0.6, 2.1, w=12.0, h=4.6)
    except FileNotFoundError:
        pass
    credit(s, "PDO from NOAA · alignment per Stier lab synthesis", y=7.1)
    speaker_note(s, "First suspect: ocean conditions. Cool, productive PDO years drove "
                 "the 1960s recovery. But we've had strings of cool, productive years "
                 "since two thousand — and the fish have not walked through the door. "
                 "Climate is necessary, but it is not what's holding them back.")

def s16_fishing():
    """S16 — Suspect 2: fishing (the average hid the extremes)"""
    s = add_slide(dark=False)
    kicker(s, "Suspect 2: Fishing", y=0.3)
    title_h1(s, "The average hid the extremes.", y=0.7, size=34)
    add_text(s, "Archipelago-wide harvest was modest. Cove-level harvest was devastating.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    try:
        img = asset("07_two_scales.png")
        add_image(s, img, 0.6, 2.1, w=8.0, h=4.6)
    except FileNotFoundError:
        pass
    # big number callouts on right
    runs = [
        {"text": "Archipelago-wide\n", "size": 13, "font": MONO, "color": C_SOFT_D},
        {"text": "4%\n\n", "size": 56, "bold": True, "font": HEAD, "color": C_KELP},
        {"text": "Cove-level\n", "size": 13, "font": MONO, "color": C_SOFT_D, "new_para": True},
        {"text": "50–70%\n\n", "size": 40, "bold": True, "font": HEAD, "color": C_RUST},
        {"text": "Local subpopulations\n", "size": 13, "font": MONO, "color": C_SOFT_D, "new_para": True},
        {"text": "up to 91%", "size": 36, "bold": True, "font": HEAD, "color": C_RUST},
    ]
    add_multi_text(s, runs, 8.9, 2.1, 4.0, 4.6)
    credit(s, "Stier et al. 2020 Ecosphere · Okamoto et al. 2020", y=7.1)
    speaker_note(s, "Second suspect: fishing. The archipelago-wide harvest rate looked "
                 "modest — about four percent. But at the cove scale, harvest rates were "
                 "fifty to seventy percent, and in some sub-populations as high as "
                 "ninety-one. The coastwide number looked safe — that was the problem. "
                 "Fishing drove the first collapse — the uncontested part of the story. "
                 "But aggregate indicators are resilience-blind. They hide the cove-scale "
                 "damage that IS the portfolio failure.")

def s17_productivity():
    """S17 — Productivity also eroded"""
    s = add_slide(dark=False)
    kicker(s, "Suspect 3: Productivity", y=0.3)
    title_h1(s, "Productivity drained, it didn't flip.", y=0.7, size=34)
    add_text(s, "Per-subpopulation process variance δ²σ fell 8-fold post-1995.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    try:
        img = asset("08_realized_growth.png")
        add_image(s, img, 0.6, 2.1, w=12.0, h=4.6)
    except FileNotFoundError:
        pass
    credit(s, "Stier et al. 2020 Ecosphere · DFO spawn-thickness data", y=7.1)
    speaker_note(s, "Third suspect: productivity. Per-subpopulation process variance "
                 "fell eight-fold between the pre-collapse period and today. Productivity "
                 "didn't flip off — it drained away over decades. This is the leading "
                 "proximate explanation for why the door didn't reopen. But there is one "
                 "more pressure that rose above all.")

def s18_mechanism():
    """S18 — Mechanism: synchrony + predator pit + lost elders/GWOF"""
    s = add_slide(dark=False)
    kicker(s, "Mechanism", y=0.3)
    title_h1(s, "Two complementary portfolio-killers.", y=0.7, size=32)
    # Two-column layout
    col_w = (DECK_W_IN - 0.8 - 0.4) / 2
    col_top = 1.7
    col_h = 5.0
    # left column: predators
    add_rect(s, 0.4, col_top, col_w, 0.5, RGBColor(0xF4, 0xF1, 0xEC))
    add_text(s, "(a) Recovered predators", 0.5, col_top + 0.05, col_w - 0.2, 0.4,
             font=HEAD, size=16, bold=True, color=C_RUST)
    runs_p = [
        {"text": "measured: ", "size": 13, "italic": True, "font": BODY, "color": C_SOFT_D},
        {"text": "synchrony rose post-1994 (φ: 0.17→0.28; Stier 2020).\n\n", "size": 13, "font": BODY},
        {"text": "hypothesis: ", "size": 13, "italic": True, "font": BODY, "color": C_SOFT_D, "new_para": True},
        {"text": "a recovered predator field is the leading explanation; ", "size": 13, "font": BODY},
        {"text": "theory predicts this signature.\n\n", "size": 13, "italic": True, "font": BODY},
        {"text": "Predator demand ≈ ⅓ of standing stock (~29% removal analogue).\n\n", "size": 13, "bold": True, "font": BODY, "new_para": True},
        {"text": "Wide-ranging predators correlate the assets — the portfolio-killer in ecological form.",
         "size": 13, "italic": True, "font": BODY, "color": C_RUST, "new_para": True},
    ]
    add_multi_text(s, runs_p, 0.5, col_top + 0.7, col_w - 0.2, col_h - 0.9)
    # right column: lost elders
    x2 = 0.4 + col_w + 0.4
    add_rect(s, x2, col_top, col_w, 0.5, RGBColor(0xF4, 0xF1, 0xEC))
    add_text(s, '(b) Lost elders / "Go With Older Fish" (GWOF)', x2 + 0.1, col_top + 0.05, col_w - 0.2, 0.4,
             font=HEAD, size=14, bold=True, color=C_RUST)
    runs_e = [
        {"text": 'Chief Gidansta (Guujaaw):\n', "size": 12, "italic": True, "font": BODY, "color": C_SOFT_D},
        {"text": '"Once herring lost the elders they lost their way to their spawning grounds."\n\n',
         "size": 14, "italic": True, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "External: MacCall et al. 2019 ICES JMS · Corten 2002 · Huse 2002/2010 · Ono et al. 2025 Nature · Jesmer et al. 2018 Science.\n\n",
         "size": 11, "font": BODY, "color": C_SOFT_D, "new_para": True},
        {"text": "Our regional age-lead test (Gear2 roe-seine, 1980–2017):\n", "size": 12, "bold": True, "font": BODY, "new_para": True},
        {"text": "ρ_firstdiff = +0.37, p ≈ 0.035, n = 32 · corroborated by mean age → spawn index at 7-yr lag (ρ = +0.43, p ≈ 0.015).\n\n",
         "size": 11, "font": MONO, "color": C_SOFT_D},
        {"text": "Consistent-with — modest, regional, descriptive. Not a demonstrated cause at sub-population scale.",
         "size": 12, "italic": True, "font": BODY, "color": C_RUST, "new_para": True},
    ]
    add_multi_text(s, runs_e, x2 + 0.1, col_top + 0.7, col_w - 0.2, col_h - 0.9)
    # bottom punch
    add_amber_rule(s, 0.4, 6.85, 0.06)
    add_text(s, "Haida knowledge and the published science name the same mechanism.",
             0.6, 6.8, 12.5, 0.4, font=HEAD, size=15, italic=True, color=C_INK_DARK)
    credit(s, "Per claim-control sheet: predator role = leading hypothesis; GWOF = consistent-with, not demonstrated at HG.", y=7.2)
    speaker_note(s, "So here's the mechanism. Two complementary portfolio-killers. "
                 "First — a recovered predator field. We measured the synchrony rise. "
                 "The leading hypothesis is that wide-ranging predators — humpbacks, "
                 "sea lions — homogenize prey by hunting everywhere at once. They "
                 "correlate the assets. Theory predicts this signature. Predator demand "
                 "today is roughly a third of the standing stock — about a twenty-nine "
                 "percent removal analogue. Second — lost elders. Chief Gidansta — "
                 "Guujaaw — said it best: 'Once herring lost the elders they lost their "
                 "way to their spawning grounds.' This is the Go-With-Older-Fish "
                 "hypothesis, and the published science has converged on it — MacCall "
                 "2019, Ono 2025, Jesmer's bighorn migration work. Our regional age-lead "
                 "test is consistent with the precursor: the ratio of repeat-to-first-"
                 "time spawners leads loss of spawning sites by about one herring "
                 "generation. Modest, regional, consistent-with — not yet demonstrated "
                 "as the cause at the subpopulation scale. The social-ecological double "
                 "helix in one slide: Haida knowledge and the published science name "
                 "the same mechanism.")

def s19_climax():
    """S19 — THE CLIMAX: P(SB<LRP)=0.38 at zero catch"""
    s = add_slide(dark=True)
    kicker(s, "The Climax", dark=True, y=0.3)
    # huge centered statistic
    add_text(s, "At zero catch since 2002,",
             0.5, 1.8, 12.3, 0.8, font=HEAD, size=32, italic=True,
             color=C_INK, align=PP_ALIGN.CENTER, dark=True)
    add_text(s, "P(SB < LRP) = 0.38",
             0.5, 2.7, 12.3, 1.5, font=HEAD, size=84, bold=True,
             color=C_RUST, align=PP_ALIGN.CENTER, dark=True)
    add_text(s, "Even at zero catch, there is a 38% chance next year's biomass falls below the Limit Reference Point.",
             0.5, 5.2, 12.3, 1.0, font=HEAD, size=18, italic=True,
             color=C_INK, align=PP_ALIGN.CENTER, dark=True)
    credit(s, "Source: DFO Pacific Herring SR 2024/2025", dark=True, y=7.1)
    speaker_note(s, "And here is where it lands. Even at zero catch since two thousand "
                 "and two, there is a thirty-eight percent chance — next year — that "
                 "spawning biomass falls below the limit reference point. We removed "
                 "the original stressor. The fish are still right at the floor. "
                 "Removing the stressor isn't enough — because the system has shifted "
                 "state.")

def s20_value_hinge():
    """S20 — VALUE HINGE: both blades of demand engine"""
    s = add_slide(dark=False)
    kicker(s, "Value Hinge", y=0.3)
    title_h1(s, "Ecology persisted. Industrial value did not.", y=0.7, size=30)
    add_text(s, "The foreign demand engine that created the BC roe fishery in 1972 has weakened on both blades.",
             0.4, 1.4, 12.5, 0.5, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    # two-blade visualization (text-led)
    col_w = (DECK_W_IN - 0.8 - 0.4) / 2
    col_top = 2.2
    col_h = 4.2
    # left blade: ecology+culture persisted
    add_text(s, "PERSISTED", 0.4, col_top, col_w, 0.4,
             font=MONO, size=12, bold=True, color=C_KELP, align=PP_ALIGN.LEFT)
    runs_left = [
        {"text": "Ecological keystone\n", "size": 16, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "Herring still the channel from plankton to predators.\n\n", "size": 13, "font": BODY},
        {"text": "Cultural keystone\n", "size": 16, "bold": True, "font": HEAD, "color": C_INK_DARK, "new_para": True},
        {"text": "kʼaaw right and ritual — never extinguished.\n", "size": 13, "font": BODY},
        {"text": '(but access has been severely constrained — see next slide)', "size": 11, "italic": True, "font": BODY, "color": C_SOFT_D},
    ]
    add_multi_text(s, runs_left, 0.4, col_top + 0.5, col_w, col_h - 0.5)
    # right blade: industrial collapsed + demand engine weakened
    x2 = 0.4 + col_w + 0.4
    add_text(s, "COLLAPSED", x2, col_top, col_w, 0.4,
             font=MONO, size=12, bold=True, color=C_RUST, align=PP_ALIGN.LEFT)
    runs_right = [
        {"text": "Industrial extractive value\n", "size": 16, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "Reduction + sac-roe → extinguished.\n\n", "size": 13, "font": BODY},
        {"text": "Foreign demand engine (Japan)\n", "size": 16, "bold": True, "font": HEAD, "color": C_INK_DARK, "new_para": True},
        {"text": "↳ Hokkaido hatchery rebuilt domestic supply (~20,000 t by 2023).\n", "size": 12, "font": BODY},
        {"text": "↳ Kazunoko demand structurally declined (demographic shift).\n\n", "size": 12, "font": BODY},
        {"text": "1980s–90s rent environment is unrecoverable on current trends.", "size": 13, "italic": True, "font": BODY, "color": C_RUST, "new_para": True},
    ]
    add_multi_text(s, runs_right, x2, col_top + 0.5, col_w, col_h - 0.5)
    # bottom punch
    add_amber_rule(s, 0.4, 6.7, 0.06)
    add_text(s, "The economic non-recovery may be more permanent than the ecological one.",
             0.6, 6.65, 12.5, 0.4, font=HEAD, size=16, italic=True, bold=True, color=C_INK_DARK)
    credit(s, "Source: 2024 HG Herring Rebuilding Plan · KCAW 2024 reporting", y=7.1)
    speaker_note(s, "Here's the value hinge — and it's the heart of why this is a "
                 "tipping-points-in-ecosystem-services talk. The ecological keystone "
                 "persisted. The cultural keystone — the Haida kʼaaw right and ritual — "
                 "persisted, though access has been severely constrained, which is the "
                 "next slide. Industrial extractive value, however, collapsed. And the "
                 "foreign demand engine that created the BC roe fishery in 1972 has "
                 "weakened on both blades — Japan rebuilt domestic supply, AND Japanese "
                 "kazunoko demand structurally declined. On current trends, the 1980s–90s "
                 "rent environment is unrecoverable. The economic non-recovery may be "
                 "more permanent than the ecological one. The right surviving is not the "
                 "same as access surviving.")

def s21_access():
    """S21 — ACCESS TIPPING POINT: 6 coupled dimensions"""
    s = add_slide(dark=False)
    kicker(s, "Access Tipping Point", y=0.3)
    title_h1(s, "The rights persisted. The access did not.", y=0.7, size=30)
    # 6-dimension grid (2 rows × 3 cols)
    dims = [
        ("(i) Regulatory", "HG commercial roe closed 2002; SOK 2004; 0 t harvest 2025; below LRP for 3 decades; DFO 150 t FSC allocation \"has not been agreed to by the Haida Nation\" (RP 2024)."),
        ("(ii) Legal", "R. v. Gladstone (1996) affirmed Heiltsuk commercial right; DFO does not recognize Haida title; 2015 Haida-led court injunction stopped a DFO opening."),
        ("(iii) Economic", "SOK $40/lb (1995) → <$6/lb (2004); HG SOK $9.1M (1988) → negligible; Davis Plan + IVQ pools = \"technological brinkmanship\" (Powell 2012)."),
        ("(iv) Physical / ecological", "Spawn extent contracted ~48% since 1940s (Gerrard 2014); Skidegate Inlet \"cleaned out … never really came back\"; Selwyn \"deserted by the larger schools.\""),
        ("(v) Infrastructure / social", "Processing plants struggle to access fish; children \"grow up without knowing the taste.\""),
        ("(vi) Cultural / kʼaaw", "\"Herring is part of the essence of who you are.\" — Barbara Wilson, 2016."),
    ]
    cell_w = (DECK_W_IN - 0.8 - 0.4) / 3
    cell_h = 1.7
    for i, (head, body) in enumerate(dims):
        row, col = divmod(i, 3)
        x = 0.4 + col * (cell_w + 0.2)
        y = 1.6 + row * (cell_h + 0.2)
        # cell bg
        add_rect(s, x, y, cell_w, cell_h, RGBColor(0xF4, 0xF1, 0xEC))
        add_text(s, head, x + 0.15, y + 0.1, cell_w - 0.3, 0.35,
                 font=HEAD, size=13, bold=True, color=C_RUST)
        add_text(s, body, x + 0.15, y + 0.5, cell_w - 0.3, cell_h - 0.6,
                 font=BODY, size=10, color=C_INK_DARK, line_spacing=1.25)
    # headline asymmetry bottom
    add_amber_rule(s, 0.4, 5.5, 0.06)
    add_text(s, "Mobile commercial fleet had access to whatever fish remained. Place-based Indigenous fishers are "
             "\"spatially constrained by boat size, fuel costs, and political/cultural boundaries\" (Stier et al. 2020) — they lost access "
             "site-by-site, decades before the regional total collapsed.",
             0.6, 5.5, 12.5, 1.4, font=HEAD, size=14, italic=True, color=C_INK_DARK, line_spacing=1.3)
    credit(s, "Sources: 2024 HG Herring Rebuilding Plan · Stier et al. 2020 Ecosphere · Powell 2012 WHQ · Gerrard 2014 · R. v. Gladstone 1996 · DFO IFMP", y=7.1)
    speaker_note(s, "And here's the second strand of the social-economic tip — access. "
                 "Because kʼaaw rights are not the same as kʼaaw access. Six coupled "
                 "dimensions, all in one lifetime. Regulatory closure. Legal — Gladstone "
                 "affirmed Heiltsuk rights but DFO does not recognize Haida commercial "
                 "title. Economic — price collapsed and licences moved out of place-"
                 "based hands. Physical — spawn extent contracted almost fifty percent. "
                 "Infrastructure decayed. And cultural — children growing up without "
                 "knowing the taste. The headline asymmetry: the mobile commercial fleet "
                 "always had access to whatever fish remained. Place-based Indigenous "
                 "fishers are spatially constrained — by boat size, fuel costs, "
                 "political and cultural boundaries — and they lost access site by site, "
                 "decades before the regional total collapsed. The rights persisted. The "
                 "fish, the licences, the infrastructure, and the access did not.")

def s22_cod_vs_herring():
    """S22 — Cod vs herring (why some recover)"""
    s = add_slide(dark=False)
    kicker(s, "Comparison", y=0.3)
    title_h1(s, "Cod rebounded. Herring didn't. Why.", y=0.7, size=34)
    add_text(s, "Stressor removed → cod recovered. Same recipe failed for HG herring.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    try:
        img = asset("sb6_cod_vs_herring.png")
        add_image(s, img, 0.6, 2.1, w=12.0, h=4.6)
    except FileNotFoundError:
        # text-only fallback
        add_text(s, "Atlantic cod, post-1992 moratorium: long delay, then recovery.\n\n"
                 "HG herring, post-2002 closure: no recovery in 22+ years.\n\n"
                 "Differences: cod lacks site-fidelity in social structure; HG herring "
                 "sits under a recovered predator field of equivalent magnitude.",
                 0.6, 2.4, 12.0, 4.0, font=HEAD, size=18, color=C_INK_DARK, line_spacing=1.5)
    credit(s, "Comparison synthesis · cod data: NAFO assessments · herring: DFO SR 2024/2025", y=7.1)
    speaker_note(s, "Quick comparison. Atlantic cod after the 1992 moratorium — closed "
                 "the fishery, and after a long delay the population has been recovering. "
                 "Same recipe — stop fishing — should work for herring. It didn't. Why? "
                 "Because cod don't carry site-fidelity in their social structure the "
                 "way herring do; and cod aren't sitting under a recovered predator "
                 "field of the same magnitude. Failed recovery here is slow drawdown "
                 "plus lost structure — not a proven new equilibrium.")

def s23_solutions():
    """S23 — SOLUTIONS: allocate the portfolio + re-couple the strands"""
    s = add_slide(dark=False)
    kicker(s, "Solutions", y=0.3)
    title_h1(s, "Allocate the portfolio. Re-couple the strands.", y=0.7, size=32)
    col_w = (DECK_W_IN - 0.8 - 0.4) / 2
    col_top = 1.7
    col_h = 5.0
    # left: Allocate
    add_text(s, "1. Allocate", 0.4, col_top, col_w, 0.4,
             font=MONO, size=13, bold=True, color=C_RUST)
    runs_l = [
        {"text": "Manage at the cove scale\n", "size": 18, "bold": True, "font": HEAD, "color": C_INK_DARK},
        {"text": "(Okamoto et al. 2020 predicts far fewer collapses with no aggregate yield loss).\n\n",
         "size": 13, "font": BODY},
        {"text": "Shift life-stage harvest\n", "size": 18, "bold": True, "font": HEAD, "color": C_INK_DARK, "new_para": True},
        {"text": "kʼaaw / SOK is non-lethal: take the eggs, leave the adults (Shelton et al. 2014).\n\n",
         "size": 13, "font": BODY},
        {"text": "Rebuild the portfolio at the scale it varies — portfolio management, literally.",
         "size": 13, "italic": True, "font": HEAD, "color": C_RUST, "new_para": True},
    ]
    add_multi_text(s, runs_l, 0.4, col_top + 0.5, col_w, col_h - 0.5)
    # right: Re-couple
    x2 = 0.4 + col_w + 0.4
    add_text(s, "2. Re-couple", x2, col_top, col_w, 0.4,
             font=MONO, size=13, bold=True, color=C_RUST)
    runs_r = [
        {"text": "Recovered predators already take ~⅓ of stock.\n", "size": 13, "font": BODY},
        {"text": "Value moved; it didn't vanish (forage + kʼaaw > roe-for-export).\n\n",
         "size": 13, "font": BODY},
        {"text": "kʼaaw — latent proof-of-concept\n", "size": 18, "bold": True, "font": HEAD, "color": C_INK_DARK, "new_para": True},
        {"text": "Right · knowledge · kin-based harvest · low-impact gear — all survive.\n", "size": 12, "font": BODY},
        {"text": "Only access is currently missing (RP 2024).\n\n", "size": 13, "italic": True, "bold": True, "font": HEAD, "color": C_RUST},
        {"text": "Allocate explicitly across ecosystem · commercial · Indigenous rights — re-coupling IS the management act. ",
         "size": 13, "font": BODY, "new_para": True},
        {"text": "Restoring access at the cove scale is the lever.", "size": 13, "italic": True, "font": HEAD, "color": C_RUST},
    ]
    add_multi_text(s, runs_r, x2, col_top + 0.5, col_w, col_h - 0.5)
    credit(s, "Sources: Okamoto et al. 2020 · Shelton et al. 2014 · 2024 HG Herring Rebuilding Plan", y=7.1)
    speaker_note(s, "Two things. First — allocate the portfolio. Manage at the cove "
                 "scale, not the archipelago scale. Okamoto and colleagues predict you "
                 "get far fewer collapses with no aggregate yield loss. And shift the "
                 "life-stage harvest — kʼaaw, spawn-on-kelp, is non-lethal; you take "
                 "the eggs and leave the adults. Second — re-couple the strands. A "
                 "recovered predator field already takes about a third of the stock. "
                 "Value moved, it didn't vanish — herring is worth more as forage and "
                 "as kʼaaw than as roe for export. And kʼaaw is the LATENT proof-of-"
                 "concept of the re-coupled regime: the right, the knowledge, the kin-"
                 "based harvest logic, the low-impact gear, all still survive. Only "
                 "access is currently missing. Re-couple them and the demonstration "
                 "becomes operational. Restoring access at the cove scale IS the "
                 "management lever.")

def s24_cogovernance():
    """S24 — Who governs it: co-governance"""
    s = add_slide(dark=False)
    kicker(s, "Who Governs", y=0.3)
    title_h1(s, "Forty years of nation-to-nation co-governance.", y=0.7, size=30)
    add_text(s, "Haida Nation · DFO · BC · Parks Canada — the 2024 Rebuilding Plan.",
             0.4, 1.4, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_SOFT_D)
    try:
        img = asset("sb5_cogovernance_timeline.png")
        add_image(s, img, 0.6, 2.1, w=12.0, h=4.4)
    except FileNotFoundError:
        # text fallback
        add_text(s, "AMB (Archipelago Management Board)\n"
                 "GayG̲ahlda Framework (Haida law) + FRRA (federal)\n"
                 "2024 Rebuilding Plan — co-developed CHN · DFO · BC · Parks Canada\n"
                 "2025 harvest recommendation: 0 t",
                 0.6, 2.4, 12.0, 4.0, font=HEAD, size=20, color=C_INK_DARK, line_spacing=1.6)
    # bottom strip
    add_amber_rule(s, 0.4, 6.55, 0.06)
    add_text(s, "Institutional milestone — not a recovery we can yet claim.",
             0.6, 6.5, 12.5, 0.4, font=HEAD, size=16, italic=True, color=C_INK_DARK)
    credit(s, "Per claim-control: co-governance = institutional milestone, not outcome metric", y=7.1)
    speaker_note(s, "And who does this? Forty years of nation-to-nation governance "
                 "precede the 2024 Rebuilding Plan. The Council of the Haida Nation. "
                 "DFO. BC. Parks Canada. The Archipelago Management Board. The "
                 "GayG̲ahlda Framework. The institution that can span both clocks "
                 "already exists. That's a milestone we can act on now — it is not a "
                 "recovery we can claim yet.")

def s25_lessons():
    """S25 — 3 transferable lessons + tipping-point reframe"""
    s = add_slide(dark=False)
    kicker(s, "Lessons", y=0.3)
    title_h1(s, "What this teaches.", y=0.7, size=34)
    # three lessons
    lessons = [
        ("1.", "Condition reference points on the current system state.",
         "SB₀ is pinned to a baseline the system has left. Below-LRP is by construction."),
        ("2.", "Match management scale to biological scale.",
         "The portfolio IS the resilience; you cannot manage it at a scale finer than the data you take."),
        ("3.", "In coupled systems, manage the ecological AND social-economic tip together.",
         "They came apart; only co-governance spans all three keystones."),
    ]
    y0 = 1.7
    row_h = 1.05
    for i, (n, head, body) in enumerate(lessons):
        y = y0 + i * row_h
        add_text(s, n, 0.4, y, 0.6, 0.5, font=HEAD, size=28, bold=True, color=C_RUST)
        add_text(s, head, 1.1, y, 11.7, 0.4, font=HEAD, size=18, bold=True, color=C_INK_DARK)
        add_text(s, body, 1.1, y + 0.4, 11.7, 0.5, font=BODY, size=13, italic=True, color=C_SOFT_D)
    # frontier hedge
    add_text(s, "Frontier (one hedge): spatial early-warning is an open hypothesis, not a result.",
             0.4, 5.1, 12.5, 0.4, font=BODY, size=12, italic=True, color=C_SOFT_D)
    # final question
    add_amber_rule(s, 0.4, 5.7, 0.06)
    add_text(s, "And — is this actually a tipping point? I'd argue no. It looks more like the slow erosion "
             "of slow variables, and the language of bifurcation may be the wrong tool here.",
             0.6, 5.65, 12.5, 1.2, font=HEAD, size=16, italic=True, color=C_INK_DARK, line_spacing=1.4)
    credit(s, "v3.6 contribution to the session: explicit tipping-point reframe", y=7.1)
    speaker_note(s, "Three transferable lessons. First: condition reference points on "
                 "the current state. SB-zero is pinned to a baseline this system has "
                 "left; below-LRP is by construction. Second: match management scale to "
                 "biological scale. The portfolio IS the resilience. Manage it at the "
                 "cove scale. Third: in coupled systems, manage the ecological and the "
                 "social-economic tip together — they came apart, and only co-governance "
                 "spans all three keystones. And one open frontier — the early-warning "
                 "signal may itself be spatial. A hypothesis from this work, not yet a "
                 "result. Finally — and this is the question I want to leave you with. "
                 "I've been calling this 'tipping points in ecosystem services' because "
                 "that's the session title. But is this actually a tipping point? I'd "
                 "argue no. It looks more like the slow erosion of slow variables. And "
                 "the language of bifurcation may be the wrong tool here.")

def s26_close():
    """S26 — Close: bookend take-home"""
    s = add_slide(dark=True)
    # same background as title
    try:
        img = asset("01_title.png")
        pic = add_image(s, img, 0, 0, DECK_W_IN, DECK_H_IN)
        spTree = pic._element.getparent()
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)
    except FileNotFoundError:
        pass
    add_rect(s, 0, 2.5, DECK_W_IN, DECK_H_IN - 2.5, C_DARK)
    kicker(s, "In one lifetime, the fish changed twice.", dark=True, y=0.3)
    add_text(s, "Recovery is a moving target.",
             0.5, 3.3, 12.3, 1.2, font=HEAD, size=58, bold=True,
             color=C_AMBER, align=PP_ALIGN.CENTER, dark=True)
    add_text(s, "Same fish. Same coast. Once it came back. Once it has not.",
             0.5, 4.7, 12.3, 0.6, font=HEAD, size=22, italic=True,
             color=C_INK, align=PP_ALIGN.CENTER, dark=True)
    add_text(s, "Thank you — to the Haida Nation, to the Pacific Herring program at DFO, and to colleagues at UCSB, OSU, UBC, and beyond.",
             0.5, 6.0, 12.3, 0.7, font=HEAD, size=14, italic=True,
             color=C_SOFT, align=PP_ALIGN.CENTER, dark=True)
    credit(s, "Adrian C. Stier · adrian.stier@ucsb.edu · Royal Society Forum · 20 May 2026",
           dark=True, y=7.05)
    speaker_note(s, "The take-home is one sentence — the same one I started with. "
                 "Recovery is a moving target. Thank you to the Haida Nation, to the "
                 "Pacific Herring program at Fisheries and Oceans Canada, and to "
                 "colleagues across UCSB, OSU, and UBC who made this work possible.")

# ─── Build sequence ─────────────────────────────────────────────────────────
def build():
    print(f"Building v3.6 deck → {OUT_PPTX}")
    s01_title()
    s02_haida_gwaii()
    s03_contrast()
    s04_herring_arrive()
    s04b_underwater()    # the underwater half of S4 reveal
    s05_wildlife()
    s06_predators()
    s07_haida()
    s08_archaeological()
    s09_three_keystones()
    s10_heterogeneity()
    s11_portfolio_theory()
    s12_portfolio_loss()
    s13_puzzle()
    s14_how_we_know()
    s15_ocean()
    s16_fishing()
    s17_productivity()
    s18_mechanism()
    s19_climax()
    s20_value_hinge()
    s21_access()
    s22_cod_vs_herring()
    s23_solutions()
    s24_cogovernance()
    s25_lessons()
    s26_close()
    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX} ({OUT_PPTX.stat().st_size/1e6:.1f} MB · {len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
